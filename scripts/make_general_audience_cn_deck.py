from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "paper" / "slides"
FIG_DIR = ROOT / "paper" / "figures"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT_DIR / "polysaccharide_kg_general_audience_cn.pptx"
SCRIPT_PATH = OUT_DIR / "polysaccharide_kg_general_audience_cn_talk_script.md"

COLORS = {
    "cream": "F7F1E5",
    "paper": "FFFDF8",
    "ink": "1F2A24",
    "muted": "667C6F",
    "green": "3A6B35",
    "sage": "CBD18F",
    "terra": "E07A5F",
    "mustard": "E3B448",
    "blue": "2F6F9F",
    "gray": "E8E2D5",
    "white": "FFFFFF",
}


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_bg(slide, color: str = "cream") -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS[color])


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None) -> None:
    if section:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(COLORS["green"])
        box.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.32), Inches(0.18), Inches(2.0), Inches(0.25))
        p = tx.text_frame.paragraphs[0]
        p.text = section
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = rgb(COLORS["green"])
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.72))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = rgb(COLORS["ink"])
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.58), Inches(1.03), Inches(11.0), Inches(0.35))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(13)
        p2.font.color.rgb = rgb(COLORS["muted"])


def add_footer(slide, idx: int) -> None:
    tx = slide.shapes.add_textbox(Inches(11.9), Inches(7.05), Inches(0.8), Inches(0.25))
    p = tx.text_frame.paragraphs[0]
    p.text = f"{idx:02d}"
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = rgb(COLORS["muted"])
    p.alignment = PP_ALIGN.RIGHT


def card(slide, x, y, w, h, title, body, color="white", accent="green") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.color.rgb = rgb(COLORS[accent])
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = rgb(COLORS[accent])
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(12)
    p2.font.color.rgb = rgb(COLORS["ink"])


def big_number(slide, x, y, number, label, color="green") -> None:
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.25), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = number
    p.font.name = "Arial"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = rgb(COLORS[color])
    p.alignment = PP_ALIGN.CENTER
    p2 = tx.text_frame.add_paragraph()
    p2.text = label
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(11)
    p2.font.color.rgb = rgb(COLORS["muted"])
    p2.alignment = PP_ALIGN.CENTER


def bullet_box(slide, x, y, w, h, bullets: list[str], fill="paper") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.color.rgb = rgb(COLORS["gray"])
    tf = shape.text_frame
    tf.margin_left = Inches(0.26)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.18)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(16)
        p.font.color.rgb = rgb(COLORS["ink"])
        p.space_after = Pt(9)


def add_image(slide, path: Path, x, y, w, h) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_arrow(slide, x1, y1, x2, y2, color="green") -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True


def add_notes(slide, notes: str) -> None:
    # python-pptx does not expose speaker notes reliably; notes are exported
    # as a separate Markdown script below.
    slide.name = notes[:30]


SLIDES: list[dict] = []


def add_slide(prs, title: str, subtitle: str | None = None, section: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle, section)
    add_footer(slide, len(prs.slides))
    return slide


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    s = add_slide(prs, "把多糖研究整理成一张“地图”", "DoLPHiN 多糖知识图谱：从零散资料到可解释检索", "开场")
    card(s, 0.7, 1.75, 4.0, 3.6, "一句话概括", "我们把天然多糖的来源、结构、功能和文献证据连成图，帮助研究者更快找到“可能有什么健康作用”。", "paper", "green")
    add_image(s, FIG_DIR / "figure1_pipeline.png", 5.0, 1.45, 7.7, 3.95)
    card(s, 7.0, 5.62, 4.0, 0.78, "今天只记住三点", "资料变地图；检索比盲目深度学习更稳；本体知识能帮一部分长尾功能。", "white", "terra")
    SLIDES.append({"title": "把多糖研究整理成一张“地图”", "notes": "大家可以把这项工作理解成：我们不是先发明一个很复杂的新模型，而是先把多糖研究里分散的信息整理成一张可以查询、可以评测的地图。今天我会用比较直观的方式讲清楚这张地图是什么、能做什么、以及目前还有什么边界。"})

    # 2
    s = add_slide(prs, "多糖是什么？为什么值得整理？", "很多天然产物的健康作用，都和复杂多糖有关", "背景")
    card(s, 0.75, 1.55, 3.5, 3.9, "生活中的多糖", "蘑菇多糖、海藻多糖、植物胶、多糖膳食纤维……它们常被研究免疫、抗氧化、代谢、凝血等作用。", "white", "green")
    card(s, 4.85, 1.55, 3.5, 3.9, "难点", "多糖不是一个小分子。它的来源、单糖组成、糖苷键、分支结构都会影响功能。", "white", "terra")
    card(s, 8.95, 1.55, 3.5, 3.9, "机会", "如果把这些信息连起来，就可能更系统地发现“结构—功能”的线索。", "white", "blue")
    SLIDES.append({"title": "多糖是什么？为什么值得整理？", "notes": "多糖并不陌生，比如真菌多糖、海藻多糖、植物来源的多糖。很多研究会报告它们有免疫调节、抗氧化、降脂、抗凝等作用。但多糖比小分子复杂得多，同一个功能可能和来源、组成、糖苷键、分支结构都有关系。所以，如果这些信息只是散落在论文里，就很难系统利用。"})

    # 3
    s = add_slide(prs, "问题不只是“数据少”，而是“证据太散”", "同一条功能线索，常常分散在不同字段和不同论文中", "问题")
    card(s, 0.8, 1.5, 2.6, 3.9, "来源", "来自哪种蘑菇、海藻或植物？", "paper", "green")
    card(s, 3.75, 1.5, 2.6, 3.9, "结构", "由哪些单糖组成？有哪些糖苷键？", "paper", "terra")
    card(s, 6.7, 1.5, 2.6, 3.9, "功能", "免疫、抗氧化、抗凝、成骨……", "paper", "blue")
    card(s, 9.65, 1.5, 2.6, 3.9, "文献", "证据来自哪篇论文？能不能追溯？", "paper", "mustard")
    add_arrow(s, 2.95, 5.8, 10.4, 5.8, "terra")
    card(s, 4.2, 6.05, 5.1, 0.55, "目标：把分散证据连成一张可用的知识地图", "", "white", "green")
    SLIDES.append({"title": "问题不只是“数据少”，而是“证据太散”", "notes": "多糖数据库里并不是完全没有信息。问题是这些信息分散在很多字段里，比如来源、组成、连接方式、功能标签和 DOI。人读一条记录可以理解，但机器很难直接用。我们的目标就是把这些字段变成节点和边，让机器也能沿着证据链检索。"})

    # 4
    s = add_slide(prs, "知识图谱：把表格变成“关系网络”", "不是只存一行记录，而是存“谁和谁有关”", "方法")
    card(s, 0.8, 1.4, 3.2, 1.1, "多糖 A", "来自某种真菌", "white", "green")
    card(s, 5.05, 1.4, 3.2, 1.1, "结构线索", "葡萄糖、半乳糖、糖苷键", "white", "terra")
    card(s, 9.3, 1.4, 3.2, 1.1, "健康功能", "免疫调节", "white", "blue")
    add_arrow(s, 3.95, 1.95, 5.0, 1.95)
    add_arrow(s, 8.25, 1.95, 9.25, 1.95)
    bullet_box(s, 1.1, 3.15, 10.9, 2.5, [
        "节点：多糖、来源生物、单糖、糖苷键、功能、疾病、论文",
        "边：来自哪里、含有什么、与什么功能相关、由哪篇论文支持",
        "好处：每个预测都可以追溯到图上的证据链"
    ])
    SLIDES.append({"title": "知识图谱：把表格变成“关系网络”", "notes": "知识图谱可以理解成关系网络。表格通常是一行一条记录，而图谱会把多糖、来源、生物结构、功能和文献都变成节点，再用边连接起来。这样一来，我们不仅知道某个多糖有什么功能，还能看到它为什么可能有这个功能，证据来自哪里。"})

    # 5
    s = add_slide(prs, "我们从 DoLPHiN 构建了一个多糖知识图谱", "规模不大，但足够做可控评测", "数据")
    nums = [("5,078", "多糖节点"), ("1,487", "来源生物"), ("23", "单糖类型"), ("323", "糖苷键"), ("66", "功能标签")]
    for i, (n, lab) in enumerate(nums):
        big_number(s, 0.75 + i * 2.45, 1.55, n, lab, ["green", "terra", "blue", "mustard", "green"][i])
    card(s, 1.15, 4.25, 10.8, 1.35, "这张图的价值", "它把“结构、来源、功能、疾病、文献”统一到一个可查询、可评测、可追溯的框架里。", "white", "green")
    SLIDES.append({"title": "我们从 DoLPHiN 构建了一个多糖知识图谱", "notes": "我们基于 DoLPHiN 数据构建了一个异构知识图谱。这里的异构意思是：图里不只有一种节点，而是有多糖、来源生物、单糖、糖苷键、功能、疾病和文献等不同类型。图的规模大约是五千多个多糖节点、六十多个功能标签。"})

    # 6
    s = add_slide(prs, "怎么判断这张图有没有用？", "把已知功能暂时遮住，看模型能不能找回来", "评测")
    card(s, 0.8, 1.5, 3.3, 3.6, "1. 遮住答案", "把一条已知“多糖—功能”关系临时隐藏。", "paper", "terra")
    card(s, 5.0, 1.5, 3.3, 3.6, "2. 让模型排序", "在 66 个功能里，看哪个功能最像正确答案。", "paper", "blue")
    card(s, 9.2, 1.5, 3.3, 3.6, "3. 看能否找回", "如果正确功能排进前 3，就说明图谱证据有帮助。", "paper", "green")
    add_arrow(s, 4.15, 3.25, 4.95, 3.25, "green")
    add_arrow(s, 8.35, 3.25, 9.15, 3.25, "green")
    SLIDES.append({"title": "怎么判断这张图有没有用？", "notes": "评测方法可以用一个考试类比。我们把已知答案暂时遮住，比如某个多糖已知有免疫调节作用，我们先把这条功能关系拿掉，然后让模型在全部功能里排序。模型如果还能把正确功能排在前面，就说明图谱里确实有可用的证据。"})

    # 7
    s = add_slide(prs, "第一个发现：简单、可解释的方法反而最稳", "当前图谱上，显式关系特征比直接 GNN 消息传递更有效", "结果")
    add_image(s, FIG_DIR / "figure2_benchmarks.png", 0.75, 1.28, 7.6, 4.4)
    card(s, 8.65, 1.55, 3.8, 1.2, "最强 clean 基线", "Meta-path + Logistic Regression\nmacro-F1 = 0.3465", "white", "green")
    card(s, 8.65, 3.15, 3.8, 1.2, "GNN 当前较弱", "Hetero GNN 约 0.046\n说明不是越复杂越好", "white", "terra")
    card(s, 8.65, 4.75, 3.8, 1.2, "解释", "有用信息集中在局部证据块里，而不是深层传播。", "white", "blue")
    SLIDES.append({"title": "第一个发现：简单、可解释的方法反而最稳", "notes": "第一个结果有点反直觉。我们试了图神经网络，但在当前这张图上，最稳的 clean 方法不是 GNN，而是把明确的关系路径拿出来，再用一个浅层分类器。这个结果并不是说 GNN 没用，而是说明当前图谱里的很多节点特征还比较弱，直接消息传递并不能自动学到好表示。"})

    # 8
    s = add_slide(prs, "第二个发现：疾病信息很强，但只能当“上界”", "它有帮助，但不能和 clean 结果混在一起解释", "结果")
    card(s, 0.85, 1.5, 3.5, 3.7, "clean 设置", "只看来源、组成、键、文献等结构证据。\n更像是在问：图谱本身有没有功能信号？", "white", "green")
    card(s, 4.9, 1.5, 3.5, 3.7, "disease-aware 设置", "加入疾病关联。\n效果更强，但疾病和功能高度相关。", "white", "terra")
    card(s, 8.95, 1.5, 3.5, 3.7, "解释边界", "这不是作弊，但必须明确：它是辅助信息上界，不是 clean 主结论。", "white", "blue")
    SLIDES.append({"title": "第二个发现：疾病信息很强，但只能当“上界”", "notes": "第二个结果是疾病信息非常强。比如某些疾病关联和功能标签本身就高度耦合，加入以后检索会更准。但这类结果不能直接当成 clean 结构功能推断，因为它回答的是另一个问题：如果我们允许使用疾病语义，最多能做到多好。"})

    # 9
    s = add_slide(prs, "第三个发现：本体知识能帮助“长尾功能”", "不是全面提升，而是帮助少见标签跨过排名门槛", "结果")
    add_image(s, FIG_DIR / "figure3_stability.png", 0.75, 1.25, 7.3, 4.25)
    card(s, 8.45, 1.45, 3.9, 1.2, "长尾问题", "有些功能样本太少，模型很难学。", "white", "terra")
    card(s, 8.45, 3.0, 3.9, 1.25, "本体的作用", "把“成骨”放到更宽的再生/组织保护语义路径中。", "white", "green")
    card(s, 8.45, 4.65, 3.9, 1.35, "稳定性", "16 个随机种子下，tail Hits@3 稳定提高。", "white", "blue")
    SLIDES.append({"title": "第三个发现：本体知识能帮助“长尾功能”", "notes": "长尾功能是指训练样本很少的功能，比如成骨这类标签。普通模型很难学，因为例子太少。我们发现本体知识不是让所有指标都大幅提高，而是在少数长尾功能上有选择性帮助。它像是给模型提供一条更宽的语义通道。"})

    # 10
    s = add_slide(prs, "一个具体例子：APP90-2 的成骨功能被救回来", "本体不是全局加分，而是窄路径上的证据转移", "案例")
    add_image(s, FIG_DIR / "figure4_case_subgraphs.png", 0.65, 1.3, 7.8, 4.45)
    card(s, 8.75, 1.5, 3.6, 1.2, "原始排序", "成骨功能排第 43", "white", "terra")
    card(s, 8.75, 3.0, 3.6, 1.2, "加入本体", "提升到第 3", "white", "green")
    card(s, 8.75, 4.5, 3.6, 1.25, "含义", "不是所有样本都变好，但少见功能能被更合理地提示出来。", "white", "blue")
    SLIDES.append({"title": "一个具体例子：APP90-2 的成骨功能被救回来", "notes": "这里是最直观的案例。APP90-2 的成骨功能在普通 disease-aware 检索里排到第 43，基本找不到。但加入 parent-child 本体传播后，提升到第 3。我们强调这不是给所有相关功能粗暴加分，而是在有足够信号时，通过窄的语义路径把少见功能提上来。"})

    # 11
    s = add_slide(prs, "也有失败：图谱还缺更细的结构语言", "失败样本告诉我们下一步该补什么", "边界")
    card(s, 0.8, 1.45, 4.0, 3.65, "CZGS-1 失败案例", "即使有来源、组成、键、疾病和 DOI，antiinflammatory 仍没有排进前列。", "white", "terra")
    card(s, 5.0, 1.45, 3.35, 3.65, "可能缺什么？", "更细的结构基序\n更准确的本体层级\n更丰富的文献证据类型", "white", "blue")
    card(s, 8.55, 1.45, 3.7, 3.65, "这很重要", "负结果不是失败，而是告诉我们图谱下一版要补哪里。", "white", "green")
    SLIDES.append({"title": "也有失败：图谱还缺更细的结构语言", "notes": "我们也保留了失败案例。CZGS-1 有不少证据，但模型还是没把抗炎排到前面。这说明当前图谱里的结构表达还不够细，比如缺少更精细的糖链结构基序，也缺少更丰富的文献证据类型。这个失败案例反而帮助我们明确下一步。"})

    # 12
    s = add_slide(prs, "这项工作真正贡献了什么？", "不是一个万能模型，而是一套可复用的图谱与评测框架", "贡献")
    bullet_box(s, 0.9, 1.45, 11.6, 4.45, [
        "把 DoLPHiN 数据转成可查询、可追溯的多糖知识图谱",
        "建立 masked poly-function retrieval 的评测任务",
        "发现当前最强 clean 信号来自可解释关系特征",
        "证明本体知识对部分长尾功能有稳定补益",
        "明确指出当前 GNN 和图谱表示的短板"
    ])
    SLIDES.append({"title": "这项工作真正贡献了什么？", "notes": "所以这项工作的贡献不是宣称发明了一个万能模型，而是构建了一套资源和评测框架。它让我们能更清楚地判断：哪些信息真的有用，哪些模型在当前图上不适合，哪些长尾功能可以通过本体知识得到帮助。"})

    # 13
    s = add_slide(prs, "谁会用到这张图？", "从数据库浏览到功能假设生成", "应用")
    card(s, 0.75, 1.55, 3.4, 3.8, "数据库使用者", "快速查看某个多糖的来源、结构、功能和文献。", "white", "green")
    card(s, 4.95, 1.55, 3.4, 3.8, "实验研究者", "优先筛选值得验证的功能假设。", "white", "terra")
    card(s, 9.15, 1.55, 3.4, 3.8, "算法研究者", "测试图谱检索、长尾学习、本体注入方法。", "white", "blue")
    SLIDES.append({"title": "谁会用到这张图？", "notes": "这张图可以服务三类人。数据库用户可以快速查证据链，实验研究者可以据此挑选更值得验证的多糖功能假设，算法研究者则可以把它当成一个图谱检索和长尾学习的测试床。"})

    # 14
    s = add_slide(prs, "总结：我们把多糖证据从“资料堆”变成“路线图”", "下一步是补更细结构、更强本体、更完整公开资源", "总结")
    card(s, 0.85, 1.45, 3.55, 3.7, "已经完成", "DoLPHiN → KG\nclean benchmark\nontology tail analysis", "white", "green")
    card(s, 4.9, 1.45, 3.55, 3.7, "主要结论", "可解释检索很强\nGNN 当前不占优\n本体帮助长尾", "white", "terra")
    card(s, 8.95, 1.45, 3.55, 3.7, "下一步", "公开图谱与代码\n补 motif-level 结构\n做更多生物验证", "white", "blue")
    SLIDES.append({"title": "总结：我们把多糖证据从“资料堆”变成“路线图”", "notes": "总结一下，我们把 DoLPHiN 的多糖记录转成知识图谱，并基于它做了可解释的功能检索评测。当前最有价值的不是复杂模型，而是图谱整理、可解释检索和长尾功能分析。下一步是把资源公开，并补充更细的结构表示和更多实验验证。"})

    # 15
    s = add_slide(prs, "谢谢，欢迎交流", "问题、建议和合作都很欢迎", "Q&A")
    card(s, 2.1, 1.8, 9.1, 2.0, "核心 takeaway", "多糖研究需要的不只是更多模型，也需要更好的证据地图。", "white", "green")
    card(s, 3.1, 4.25, 7.1, 1.25, "联系信息", "作者 / 单位 / 邮箱：提交前替换", "paper", "terra")
    SLIDES.append({"title": "谢谢，欢迎交流", "notes": "最后一句话总结：多糖功能研究需要的不只是更复杂的模型，也需要更好的证据地图。谢谢大家，欢迎提问和交流。"})

    # backups
    s = add_slide(prs, "备份：如果有人问“为什么不用更复杂的 GNN？”", "我们的回答：不是不用，而是当前图谱表示还不支持它发挥优势", "备份")
    bullet_box(s, 1.0, 1.5, 11.2, 4.5, [
        "GNN 和 no-message ablation 表现接近",
        "说明消息传递没有提供额外有效信息",
        "可能原因：非多糖节点特征太弱、结构 motif 表达不足",
        "未来方向：先增强图谱表示，再谈更强 GNN"
    ])
    SLIDES.append({"title": "备份：为什么不用更复杂的 GNN？", "notes": "如果有人问为什么不用更复杂的 GNN，可以回答：我们用了，也做了 no-message ablation。结果说明问题不主要在模型深度，而在当前图谱表示还不够强。"})

    s = add_slide(prs, "备份：clean 与 disease-aware 有什么区别？", "这是审稿人最可能问的问题之一", "备份")
    bullet_box(s, 1.0, 1.5, 11.2, 4.6, [
        "clean：不用 disease 信息，检验结构/来源/文献证据本身",
        "disease-aware：允许疾病语义，作为辅助信息上界",
        "两者不能混报，否则会高估结构—功能推断能力",
        "我们已修复 clean poly_x 中的 disease-degree 泄漏"
    ])
    SLIDES.append({"title": "备份：clean 与 disease-aware 有什么区别？", "notes": "clean 和 disease-aware 的区别非常重要。clean 是主科学结论，disease-aware 是辅助上界。我们也修复了 clean 特征中的 disease-degree 泄漏，确保结果边界清楚。"})

    s = add_slide(prs, "备份：下一版图谱最应该补什么？", "让模型看到更接近生物机制的证据", "备份")
    bullet_box(s, 1.0, 1.5, 11.2, 4.6, [
        "更细的糖链 motif / branching 表达",
        "更标准的功能本体和疾病本体映射",
        "文献证据类型：实验条件、剂量、模型体系",
        "更严格的 source-level / publication-level split"
    ])
    SLIDES.append({"title": "备份：下一版图谱最应该补什么？", "notes": "下一版最值得补的是更接近生物机制的证据，比如糖链 motif、分支结构、实验条件、剂量、模型体系等。这样才能让模型不只是看标签共现，而是看到更真实的结构功能关系。"})

    prs.save(PPTX_PATH)


def write_script() -> None:
    lines = [
        "# 中文讲稿：DoLPHiN 多糖知识图谱（普通听众版）",
        "",
        "建议时长：12-15 分钟。主讲 15 页，备份 3 页仅用于问答。",
        "",
    ]
    for i, item in enumerate(SLIDES, start=1):
        lines.append(f"## Slide {i}. {item['title']}")
        lines.append("")
        lines.append(item["notes"])
        lines.append("")
    SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    build_deck()
    write_script()
    print({"pptx": str(PPTX_PATH), "script": str(SCRIPT_PATH)})

