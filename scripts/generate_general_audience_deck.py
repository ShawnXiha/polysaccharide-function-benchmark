"""Generate a Chinese general-audience PPT deck and speaker script."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentations"
PPTX_PATH = OUT_DIR / "polysaccharide_function_benchmark_general_audience_zh.pptx"
SCRIPT_PATH = OUT_DIR / "polysaccharide_function_benchmark_general_audience_zh_script.md"


WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(39, 48, 51),
    "muted": RGBColor(91, 105, 112),
    "tea": RGBColor(63, 122, 87),
    "tea_dark": RGBColor(33, 89, 63),
    "cream": RGBColor(248, 244, 232),
    "sand": RGBColor(232, 219, 190),
    "gold": RGBColor(205, 154, 74),
    "coral": RGBColor(216, 104, 78),
    "blue": RGBColor(70, 119, 153),
    "white": RGBColor(255, 255, 255),
    "gray": RGBColor(225, 228, 226),
}

FONT = "Microsoft YaHei"


def add_textbox(slide, x, y, w, h, text, size=24, color="ink", bold=False, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    if align is not None:
        p.alignment = align
    return box


def add_title(slide, text, subtitle=None):
    add_textbox(slide, Inches(0.65), Inches(0.35), Inches(11.8), Inches(0.55), text, 30, "ink", True)
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(0.92), Inches(11.4), Inches(0.32), subtitle, 13, "muted")
    add_rule(slide, Inches(0.65), Inches(1.25), Inches(2.4), "tea")


def add_rule(slide, x, y, w, color="tea"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS[color]
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="tea", icon=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["gray"]
    card.line.width = Pt(1)
    add_rule(slide, x + Inches(0.18), y + Inches(0.2), Inches(0.72), accent)
    if icon:
        add_textbox(slide, x + Inches(0.18), y + Inches(0.38), Inches(0.55), Inches(0.45), icon, 21, accent, True)
        title_x = x + Inches(0.8)
        title_w = w - Inches(1.0)
    else:
        title_x = x + Inches(0.22)
        title_w = w - Inches(0.44)
    add_textbox(slide, title_x, y + Inches(0.35), title_w, Inches(0.35), title, 18, "ink", True)
    body_box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.9), w - Inches(0.44), h - Inches(1.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(body):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(15)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(5)
    return card


def add_footer(slide, num):
    add_textbox(slide, Inches(11.9), Inches(7.05), Inches(0.8), Inches(0.25), f"{num:02d}", 10, "muted", align=PP_ALIGN.RIGHT)


def add_background(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["cream"]
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.95), WIDE_W, Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(238, 232, 215)
    band.line.fill.background()


def add_big_number(slide, x, y, number, label, color="tea"):
    add_textbox(slide, x, y, Inches(2.3), Inches(0.75), number, 38, color, True)
    add_textbox(slide, x, y + Inches(0.72), Inches(2.4), Inches(0.45), label, 13, "muted")


def add_bar(slide, x, y, label, value, max_value, color="tea", note=None):
    add_textbox(slide, x, y, Inches(2.0), Inches(0.3), label, 13, "ink", True)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(2.2), y + Inches(0.05), Inches(3.5), Inches(0.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["gray"]
    bg.line.fill.background()
    fg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + Inches(2.2),
        y + Inches(0.05),
        Inches(3.5 * value / max_value),
        Inches(0.2),
    )
    fg.fill.solid()
    fg.fill.fore_color.rgb = COLORS[color]
    fg.line.fill.background()
    add_textbox(slide, x + Inches(5.85), y - Inches(0.02), Inches(0.85), Inches(0.3), f"{value:.3f}", 12, color, True)
    if note:
        add_textbox(slide, x + Inches(6.75), y - Inches(0.02), Inches(2.2), Inches(0.3), note, 11, "muted")


def add_chip(slide, x, y, text, color="tea"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.fill.background()
    add_textbox(slide, x + Inches(0.05), y + Inches(0.08), Inches(1.9), Inches(0.25), text, 12, "white", True, PP_ALIGN.CENTER)


def make_slide(prs, index, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, index)
    return slide


def create_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.75), Inches(0.75), Inches(11.6), Inches(1.0), "天然多糖能不能用机器学习预测功能？", 34, "ink", True)
    add_textbox(slide, Inches(0.8), Inches(1.65), Inches(10.8), Inches(0.45), "一个面向公开数据的结构-功能预测基准研究", 20, "tea_dark")
    add_card(
        slide,
        Inches(0.85),
        Inches(2.55),
        Inches(3.6),
        Inches(2.4),
        "一句话结论",
        ["现在最可靠的发现不是“模型很强”", "而是：分子量、单糖残基、支链信息最有用"],
        "gold",
        "✓",
    )
    add_card(
        slide,
        Inches(4.85),
        Inches(2.55),
        Inches(3.6),
        Inches(2.4),
        "研究对象",
        ["天然多糖：食品、植物、菌菇、海藻等来源", "目标：预测已报道的生物功能"],
        "tea",
        "◎",
    )
    add_card(
        slide,
        Inches(8.85),
        Inches(2.55),
        Inches(3.6),
        Inches(2.4),
        "适合听众",
        ["不需要懂深度学习", "只需要理解：数据、特征、验证"],
        "blue",
        "→",
    )
    add_footer(slide, 1)

    # 2
    slide = make_slide(prs, 2, "为什么关心天然多糖？", "它们常被研究为食品与健康功能相关成分")
    add_card(slide, Inches(0.8), Inches(1.65), Inches(3.7), Inches(2.1), "常见功能", ["抗氧化", "免疫调节", "抗肿瘤 / 代谢调节"], "tea", "1")
    add_card(slide, Inches(4.8), Inches(1.65), Inches(3.7), Inches(2.1), "现实问题", ["文献很多，但信息分散", "结构描述不统一"], "gold", "2")
    add_card(slide, Inches(8.8), Inches(1.65), Inches(3.7), Inches(2.1), "机器学习机会", ["把公开记录整理成表", "看看哪些结构信号真的有用"], "blue", "3")
    add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.1), Inches(0.8), "不是要替代实验，而是帮助研究者更快地筛选候选多糖。", 24, "tea_dark", True, PP_ALIGN.CENTER)

    # 3
    slide = make_slide(prs, 3, "难点：多糖不像蛋白质那样有一条清晰序列", "公开数据库里的记录往往是不完整、混合、文本化的")
    add_card(slide, Inches(0.75), Inches(1.55), Inches(3.7), Inches(2.7), "蛋白质/小分子", ["通常有较标准的序列或结构表示", "适合直接喂给成熟模型"], "blue", "A")
    add_card(slide, Inches(4.8), Inches(1.55), Inches(3.7), Inches(2.7), "天然多糖", ["可能只有组成、支链、分子量范围", "键型和修饰信息经常缺失"], "coral", "B")
    add_card(slide, Inches(8.85), Inches(1.55), Inches(3.7), Inches(2.7), "我们的选择", ["先不追求复杂模型", "先问：哪些公开信号最稳定？"], "tea", "C")
    add_textbox(slide, Inches(1.0), Inches(5.05), Inches(11.2), Inches(0.5), "核心思路：把问题先做小、做清楚、做可复现。", 24, "ink", True, PP_ALIGN.CENTER)

    # 4
    slide = make_slide(prs, 4, "我们做了什么：从网页记录到可复现实验", "这项工作更像“搭建基准 + 检查结构信号”，不是炫技模型")
    steps = [
        ("抓取", "DoLPHiN 公开记录"),
        ("清洗", "统一字段和标签"),
        ("建表", "4121 条记录 / 18 类功能"),
        ("建模", "比较多种简单基线"),
        ("验证", "随机切分 + DOI 分组切分"),
    ]
    x0 = 0.65
    for i, (head, body) in enumerate(steps):
        x = Inches(x0 + i * 2.45)
        add_card(slide, x, Inches(2.0), Inches(2.05), Inches(2.0), head, [body], ["tea", "gold", "blue", "coral", "tea"][i], str(i + 1))
        if i < len(steps) - 1:
            add_textbox(slide, Inches(x0 + 2.12 + i * 2.45), Inches(2.75), Inches(0.35), Inches(0.4), "→", 22, "muted", True)
    add_textbox(slide, Inches(0.9), Inches(5.25), Inches(11.4), Inches(0.45), "输出不是一个神奇黑箱，而是一套别人可以复查的数据和实验流程。", 22, "tea_dark", True, PP_ALIGN.CENTER)

    # 5
    slide = make_slide(prs, 5, "这个数据集长什么样？", "公开数据可用，但有明显弱监督和缺失问题")
    add_big_number(slide, Inches(1.0), Inches(1.65), "4,121", "保留记录", "tea")
    add_big_number(slide, Inches(4.0), Inches(1.65), "18", "功能标签", "blue")
    add_big_number(slide, Inches(7.0), Inches(1.65), "1.56", "平均每条记录标签数", "gold")
    add_big_number(slide, Inches(10.1), Inches(1.65), "80.9%", "键型信息缺失", "coral")
    add_card(slide, Inches(1.0), Inches(4.15), Inches(5.3), Inches(1.55), "重要提醒", ["没有某个功能标签，不代表一定没有这个功能", "更可能是“文献还没有报道”"], "coral", "!")
    add_card(slide, Inches(7.0), Inches(4.15), Inches(5.3), Inches(1.55), "因此我们怎么解释结果", ["把指标当作相对比较", "不把分数解释成绝对生物真相"], "tea", "i")

    # 6
    slide = make_slide(prs, 6, "我们比较了哪些方法？", "从最简单的频率预测，到文本模型、图模型和结构化特征")
    add_chip(slide, Inches(0.95), Inches(1.65), "多数类预测", "gray")
    add_chip(slide, Inches(3.0), Inches(1.65), "逻辑回归", "blue")
    add_chip(slide, Inches(5.05), Inches(1.65), "随机森林", "tea")
    add_chip(slide, Inches(7.1), Inches(1.65), "SVM / TF-IDF", "gold")
    add_chip(slide, Inches(9.15), Inches(1.65), "序列模型", "coral")
    add_chip(slide, Inches(11.2), Inches(1.65), "图模型", "tea")
    add_card(slide, Inches(1.0), Inches(3.1), Inches(5.2), Inches(2.1), "最强的实用锚点", ["不是复杂神经网络", "而是调好参数的稀疏逻辑回归"], "blue", "★")
    add_card(slide, Inches(7.0), Inches(3.1), Inches(5.2), Inches(2.1), "我们的结构化增强", ["只加入三类简单信息", "分子量 + 残基集合 + 支链"], "tea", "＋")

    # 7
    slide = make_slide(prs, 7, "主结果：提升很小，而且不够稳健", "这也是论文现在最诚实、最可信的地方")
    add_bar(slide, Inches(1.0), Inches(1.8), "多数类", 0.0397, 0.30, "gray")
    add_bar(slide, Inches(1.0), Inches(2.35), "未调逻辑回归", 0.1580, 0.30, "blue")
    add_bar(slide, Inches(1.0), Inches(2.9), "Graph GCN", 0.2103, 0.30, "tea", "已修 vocab 泄漏")
    add_bar(slide, Inches(1.0), Inches(3.45), "调参逻辑回归", 0.2610, 0.30, "gold", "强锚点")
    add_bar(slide, Inches(1.0), Inches(4.0), "poly-core v1", 0.2678, 0.30, "coral", "小幅提升")
    add_card(slide, Inches(8.0), Inches(1.7), Inches(4.2), Inches(3.0), "关键解释", ["随机切分下有小幅提升", "bootstrap 置信区间跨 0", "DOI 分组切分下优势消失"], "coral", "!")
    add_textbox(slide, Inches(1.0), Inches(5.65), Inches(11.2), Inches(0.45), "所以我们不说“新模型赢了”，而说“发现了哪些结构信号更值得相信”。", 21, "ink", True, PP_ALIGN.CENTER)

    # 8
    slide = make_slide(prs, 8, "真正有用的是三个结构信号", "消融实验告诉我们：不是所有看起来合理的特征都有帮助")
    add_card(slide, Inches(0.9), Inches(1.65), Inches(3.7), Inches(2.2), "分子量 MW", ["贡献最大", "可能反映聚合度、溶解性、活性差异"], "tea", "1")
    add_card(slide, Inches(4.85), Inches(1.65), Inches(3.7), Inches(2.2), "残基集合", ["例如葡萄糖、半乳糖、阿拉伯糖等", "帮助区分结构组成"], "blue", "2")
    add_card(slide, Inches(8.8), Inches(1.65), Inches(3.7), Inches(2.2), "支链信息", ["作用较小但仍有帮助", "说明分支结构值得保留"], "gold", "3")
    add_card(slide, Inches(2.0), Inches(4.55), Inches(4.3), Inches(1.35), "不太有效", ["修饰字段、来源王国、粗略组成计数"], "coral", "×")
    add_card(slide, Inches(7.0), Inches(4.55), Inches(4.3), Inches(1.35), "研究启发", ["先相信稳定简单信号，再谈复杂模型"], "tea", "✓")

    # 9
    slide = make_slide(prs, 9, "我们也发现了失败：严格切分后方法会退化", "这说明公开数据里存在论文来源和记录风格的影响")
    add_card(slide, Inches(0.9), Inches(1.55), Inches(5.5), Inches(2.2), "随机切分", ["训练集和测试集可能来自相近论文", "容易学到记录风格"], "gold", "A")
    add_card(slide, Inches(6.95), Inches(1.55), Inches(5.5), Inches(2.2), "DOI 分组切分", ["同一论文的记录不会跨集合", "更接近真实泛化"], "blue", "B")
    add_textbox(slide, Inches(1.1), Inches(4.45), Inches(11.0), Inches(0.5), "结果：poly-core 在随机切分略好，但在 DOI 分组切分不再优于强基线。", 22, "coral", True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.1), Inches(5.25), Inches(11.0), Inches(0.4), "这不是失败，而是帮我们划清了当前证据边界。", 20, "tea_dark", True, PP_ALIGN.CENTER)

    # 10
    slide = make_slide(prs, 10, "这项工作的价值在哪里？", "更像是一张地图，而不是一个终点答案")
    add_card(slide, Inches(0.9), Inches(1.55), Inches(3.6), Inches(2.4), "给研究者", ["知道哪些字段值得优先整理", "减少盲目堆模型"], "tea", "人")
    add_card(slide, Inches(4.85), Inches(1.55), Inches(3.6), Inches(2.4), "给数据库建设", ["分子量、残基、支链要更标准", "功能证据类型要补齐"], "blue", "库")
    add_card(slide, Inches(8.8), Inches(1.55), Inches(3.6), Inches(2.4), "给机器学习", ["天然多糖不是标准 glycan 任务", "需要弱监督和分组验证"], "gold", "ML")
    add_textbox(slide, Inches(1.0), Inches(5.25), Inches(11.4), Inches(0.45), "下一步不是“更大的模型”，而是“更干净的数据 + 更严格的验证”。", 22, "ink", True, PP_ALIGN.CENTER)

    # 11
    slide = make_slide(prs, 11, "下一步计划", "把现在的 benchmark 变成更可靠的研究工具")
    add_card(slide, Inches(0.95), Inches(1.5), Inches(3.7), Inches(2.5), "数据侧", ["补 evidence type", "规范分子量单位", "改善跨库标签对齐"], "tea", "1")
    add_card(slide, Inches(4.85), Inches(1.5), Inches(3.7), Inches(2.5), "方法侧", ["PU 学习", "label ranking", "更稳健的 grouped split"], "blue", "2")
    add_card(slide, Inches(8.75), Inches(1.5), Inches(3.7), Inches(2.5), "应用侧", ["候选多糖筛选", "实验优先级排序", "数据库质量反馈"], "gold", "3")
    add_textbox(slide, Inches(1.0), Inches(5.25), Inches(11.2), Inches(0.45), "目标：让模型成为实验和数据库建设的辅助工具，而不是替代专家判断。", 21, "tea_dark", True, PP_ALIGN.CENTER)

    # 12
    slide = make_slide(prs, 12, "总结：这篇工作的主张很克制", "三个可以带走的结论")
    add_card(slide, Inches(1.0), Inches(1.55), Inches(3.4), Inches(2.5), "问题重要", ["天然多糖功能预测有现实意义", "但公开数据仍不完美"], "tea", "1")
    add_card(slide, Inches(4.95), Inches(1.55), Inches(3.4), Inches(2.5), "强基线很强", ["简单稀疏模型不应被低估", "复杂模型未必稳定"], "blue", "2")
    add_card(slide, Inches(8.9), Inches(1.55), Inches(3.4), Inches(2.5), "信号很清楚", ["分子量、残基、支链最值得保留", "其他特征要谨慎"], "gold", "3")
    add_textbox(slide, Inches(1.2), Inches(5.4), Inches(10.8), Inches(0.45), "谢谢，欢迎交流。", 28, "ink", True, PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


SCRIPT = """# 中文通俗版讲稿

## 01. 天然多糖能不能用机器学习预测功能？

今天我想用比较通俗的方式介绍这项工作。我们研究的是天然多糖的结构和功能之间有没有可以被机器学习利用的规律。这里的重点不是证明某个模型特别强，而是想回答一个更基础的问题：在当前公开数据条件下，哪些结构信息真的有用。

一句话总结：现在最可靠的发现不是“模型赢了”，而是分子量、单糖残基集合、支链信息这三类结构信号最值得保留。

## 02. 为什么关心天然多糖？

天然多糖广泛存在于植物、菌菇、海藻和食品原料里。很多文献会研究它们的抗氧化、免疫调节、抗肿瘤、代谢调节等功能。

现实问题是，相关文献很多，但信息很分散。不同论文对结构的写法也不统一。所以我们希望用机器学习把这些公开记录整理起来，帮助研究者更快地筛选候选多糖。

这里要强调：模型不是替代实验，而是辅助实验设计。

## 03. 难点：多糖不像蛋白质那样有一条清晰序列

蛋白质通常有比较清楚的氨基酸序列，小分子也有比较标准的结构表示。但天然多糖不一样。公开记录里经常只有单糖组成、分子量范围、支链描述，有时候键型和修饰信息还缺失。

所以如果直接套用成熟的 glycan 模型，效果未必可靠。我们的策略是先不追求复杂模型，而是把问题收缩到：哪些公开可获得的结构信号最稳定。

## 04. 我们做了什么？

整个流程可以理解成五步。

第一，从 DoLPHiN 公网站点抓取天然多糖记录。第二，对字段和功能标签做统一。第三，构建一个可监督学习的数据表。第四，比较多种简单和复杂模型。第五，用随机切分和 DOI 分组切分做验证。

这样做的价值是：别人可以复查数据、复查代码、复查实验流程。

## 05. 数据集长什么样？

最终主数据集有 4121 条记录，保留 18 类功能标签。平均每条记录大约 1.56 个标签，说明这是一个比较稀疏的多标签任务。

同时，数据也不完美。例如键型信息缺失率很高，修饰字段几乎不可用。更重要的是，文献里没有报道某个功能，不代表这个多糖真的没有这个功能，只能说明目前没有记录。

所以我们把这个任务看成弱监督任务，而不是完全可靠的真负样本分类。

## 06. 我们比较了哪些方法？

我们比较了从简单到复杂的一系列方法，包括多数类预测、逻辑回归、随机森林、TF-IDF、SVM、序列模型、图模型。

最后发现，最值得作为锚点的不是复杂神经网络，而是调好参数的稀疏逻辑回归。我们的 poly-core v1 也是在这个强基线基础上，只加入三类结构化信号：分子量、残基集合和支链。

## 07. 主结果：提升很小，而且不够稳健

从数字上看，随机切分下 poly-core v1 的 macro-F1 从 0.2610 提高到 0.2678，确实有一点提升。

但是这个提升很小，bootstrap 置信区间跨过 0；在更严格的 DOI 分组切分下，这个优势还会消失。

所以我们不说“新模型显著胜出”，而说“结构化特征提供了一些信号，但还不够稳健”。

## 08. 真正有用的是三个结构信号

消融实验告诉我们，最有用的是三个信号。

第一是分子量。去掉分子量后性能下降最大。第二是单糖残基集合，例如葡萄糖、半乳糖、阿拉伯糖等。第三是支链信息，贡献较小但仍然有帮助。

相反，一些看起来合理的特征，例如修饰字段、来源王国、粗略组成计数，并没有稳定帮助。

## 09. 严格切分后为什么会退化？

随机切分有一个风险：同一篇论文或相近论文里的记录可能同时出现在训练集和测试集，这会让模型学到论文写作风格，而不是结构和功能的真实关系。

所以我们增加了 DOI 分组切分，同一篇论文的记录不会跨集合。结果显示，poly-core 的优势不再稳定。

这不是坏事。它帮助我们划清证据边界，避免把随机切分下的小优势说得太大。

## 10. 这项工作的价值在哪里？

我认为这项工作的价值主要有三点。

第一，给多糖研究者一个可复现的 benchmark。第二，告诉数据库建设者哪些字段值得优先标准化。第三，提醒机器学习研究者：天然多糖不是标准 glycan 任务，需要弱监督和更严格的分组验证。

换句话说，这更像是一张地图，而不是一个终点答案。

## 11. 下一步计划

下一步可以从三个方向推进。

数据侧，需要补充 evidence type，规范分子量单位，改善跨数据库标签对齐。方法侧，可以尝试 positive-unlabeled learning、label ranking 和更稳定的 grouped split。应用侧，可以把模型用于候选多糖筛选和实验优先级排序。

目标不是让模型替代专家，而是让模型帮助专家更快定位值得实验验证的对象。

## 12. 总结

最后总结三句话。

第一，天然多糖功能预测是重要但困难的问题。第二，简单强基线非常重要，复杂模型并不天然可靠。第三，在当前公开数据条件下，分子量、残基集合和支链信息是最清楚、最值得保留的结构信号。

谢谢，欢迎交流。
"""


def write_script() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    SCRIPT_PATH.write_text(SCRIPT, encoding="utf-8")


def main() -> None:
    create_deck()
    write_script()
    print(f"Wrote {PPTX_PATH}")
    print(f"Wrote {SCRIPT_PATH}")


if __name__ == "__main__":
    main()
